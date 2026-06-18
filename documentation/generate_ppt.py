from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color Palette ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)   # Dark Navy (background utama)
BLUE      = RGBColor(0x1A, 0x5F, 0x7A)   # Medium Blue
ACCENT    = RGBColor(0x00, 0xC2, 0xFF)   # Cyan Accent
YELLOW    = RGBColor(0xFF, 0xC8, 0x00)   # Yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE8, 0xF4, 0xFD)   # Very Light Blue
GRAY      = RGBColor(0xA0, 0xB4, 0xC8)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
ORANGE    = RGBColor(0xF3, 0x96, 0x14)
RED       = RGBColor(0xE7, 0x4C, 0x3C)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, text, l, t, w, h, font_size=14, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    # Accent top bar
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)

    # Left decorative rect
    add_rect(slide, 0, 0.08, 0.5, 7.42, BLUE)

    # Big title background block
    add_rect(slide, 0.5, 1.5, 9, 2.2, BLUE)

    # Accent corner
    add_rect(slide, 9.5, 1.5, 0.08, 2.2, YELLOW)

    # LIVO brand
    add_textbox(slide, "LIVO", 0.8, 0.2, 6, 1.1, font_size=54, bold=True, color=ACCENT)
    add_textbox(slide, "Learning Innovation", 0.8, 1.05, 6, 0.5, font_size=16, color=YELLOW, bold=False)

    # Main title
    add_textbox(slide, "Product Knowledge", 0.8, 1.6, 8.5, 0.7, font_size=30, bold=True, color=WHITE)
    add_textbox(slide, "Sistem Manajemen Bimbingan Belajar", 0.8, 2.3, 8.5, 0.5, font_size=18, color=LIGHT)

    # Divider
    add_rect(slide, 0.8, 3.0, 5, 0.04, ACCENT)

    # Subtitle info
    add_textbox(slide, "Platform digital terpadu untuk pengelolaan operasional bimbel\nsecara efisien & terstruktur", 0.8, 3.15, 9, 0.8, font_size=13, color=GRAY, italic=True)

    # Bottom info
    add_rect(slide, 0, 6.9, 13.33, 0.6, BLUE)
    add_textbox(slide, "Srengseng Sawah, Jakarta Selatan  |  Matematika & Bahasa Inggris  |  TK – SMP",
                0.5, 6.95, 12, 0.45, font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)

    # Right visual elements
    add_rect(slide, 10.5, 1.8, 2, 1.6, RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(slide, "📚", 10.9, 1.85, 1.2, 0.8, font_size=40, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Bimbel Digital", 10.5, 2.7, 2, 0.4, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    return slide


def add_section_divider(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)
    add_rect(slide, 0.18, 0, 13.15, 7.5, BLUE)

    # Number/label
    add_rect(slide, 1, 2.5, 11, 0.08, YELLOW)

    add_textbox(slide, title, 1, 2.8, 11, 1.2, font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle, 1, 4.0, 11, 0.6, font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)

    return slide


def add_content_slide(prs, title, bullets, icon="▶", two_col=False, col2_bullets=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    # Header bar
    add_rect(slide, 0, 0, 13.33, 1.1, BLUE)
    add_rect(slide, 0, 0, 0.06, 1.1, ACCENT)
    add_textbox(slide, f"{icon}  {title}", 0.2, 0.15, 12, 0.75, font_size=22, bold=True, color=WHITE)

    # Bottom accent
    add_rect(slide, 0, 7.0, 13.33, 0.5, RGBColor(0x0A, 0x14, 0x20))
    add_textbox(slide, "LIVO — Learning Innovation", 0, 7.05, 13.33, 0.35, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    if not two_col:
        # Single column bullets
        y = 1.25
        for bullet in bullets:
            if bullet.startswith("##"):
                # Sub-header
                add_rect(slide, 0.3, y, 12.5, 0.35, RGBColor(0x1A, 0x3A, 0x5A))
                add_textbox(slide, bullet[2:].strip(), 0.5, y + 0.03, 12, 0.28, font_size=13, bold=True, color=ACCENT)
                y += 0.48
            else:
                # Normal bullet
                add_textbox(slide, "◆", 0.3, y, 0.3, 0.35, font_size=9, color=YELLOW)
                add_textbox(slide, bullet, 0.65, y, 12, 0.38, font_size=12.5, color=WHITE)
                y += 0.42
            if y > 6.7:
                break
    else:
        # Two column layout
        col1_w = 5.8
        col2_x = 6.9

        add_textbox(slide, bullets[0] if bullets else "Kolom 1", 0.3, 1.2, col1_w, 0.4, font_size=13, bold=True, color=ACCENT)
        add_rect(slide, 0.3, 1.6, col1_w, 0.04, ACCENT)
        y = 1.75
        for b in bullets[1:]:
            add_textbox(slide, "◆  " + b, 0.3, y, col1_w, 0.38, font_size=12, color=WHITE)
            y += 0.42

        if col2_bullets:
            add_textbox(slide, col2_bullets[0], col2_x, 1.2, col1_w, 0.4, font_size=13, bold=True, color=YELLOW)
            add_rect(slide, col2_x, 1.6, col1_w, 0.04, YELLOW)
            y2 = 1.75
            for b in col2_bullets[1:]:
                add_textbox(slide, "◆  " + b, col2_x, y2, col1_w, 0.38, font_size=12, color=WHITE)
                y2 += 0.42

    return slide


def add_flow_slide(prs, title, steps):
    """Slide with numbered flow steps."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, 13.33, 1.0, BLUE)
    add_rect(slide, 0, 0, 0.06, 1.0, ACCENT)
    add_textbox(slide, title, 0.2, 0.13, 12, 0.72, font_size=22, bold=True, color=WHITE)

    add_rect(slide, 7.0, 0, 0.04, 1.0, YELLOW)

    colors = [ACCENT, YELLOW, GREEN, ORANGE, RGBColor(0xAB, 0x47, 0xBC), RED]
    per_row = 3
    box_w = 3.8
    box_h = 1.6
    x_start = 0.3
    y_start = 1.15
    gap_x = 0.25
    gap_y = 0.3

    for i, (step_title, step_desc) in enumerate(steps):
        row = i // per_row
        col = i % per_row
        x = x_start + col * (box_w + gap_x)
        y = y_start + row * (box_h + gap_y)
        c = colors[i % len(colors)]

        # Card bg
        add_rect(slide, x, y, box_w, box_h, RGBColor(0x12, 0x28, 0x40))

        # Number badge
        add_rect(slide, x, y, 0.5, box_h, c)
        add_textbox(slide, str(i + 1), x, y + 0.5, 0.5, 0.5, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Content
        add_textbox(slide, step_title, x + 0.55, y + 0.1, box_w - 0.65, 0.4, font_size=12, bold=True, color=c)
        add_textbox(slide, step_desc, x + 0.55, y + 0.52, box_w - 0.65, 0.95, font_size=10.5, color=LIGHT)

    add_rect(slide, 0, 7.0, 13.33, 0.5, RGBColor(0x0A, 0x14, 0x20))
    add_textbox(slide, "LIVO — Learning Innovation", 0, 7.05, 13.33, 0.35, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    return slide


def add_role_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, 13.33, 1.0, BLUE)
    add_rect(slide, 0, 0, 0.06, 1.0, ACCENT)
    add_textbox(slide, "👥  Pengguna Sistem LIVO", 0.2, 0.13, 12, 0.72, font_size=22, bold=True, color=WHITE)

    roles = [
        ("🏫", "Admin", ACCENT,
         ["Mengelola seluruh data sistem", "Menyetujui/menolak pembayaran", "Membuat jadwal kelas", "Mempublikasi laporan evaluasi", "Mengelola pendaftaran siswa"]),
        ("👨‍🏫", "Tutor", YELLOW,
         ["Melihat jadwal mengajar", "Input absensi siswa", "Input nilai tugas & kuis", "Menulis catatan perkembangan", "Memperbarui profil & keahlian"]),
        ("👨‍👩‍👧", "Orang Tua / Siswa", GREEN,
         ["Daftar secara online", "Upload bukti pembayaran", "Melihat laporan nilai", "Melihat catatan performa", "Konsultasi tugas 24/7"]),
    ]

    x_positions = [0.3, 4.6, 8.9]
    card_w = 4.0

    for i, (icon, role, color, items) in enumerate(roles):
        x = x_positions[i]
        # Card
        add_rect(slide, x, 1.1, card_w, 5.6, RGBColor(0x0E, 0x20, 0x35))
        add_rect(slide, x, 1.1, card_w, 0.08, color)

        # Icon circle bg
        add_rect(slide, x + 1.5, 1.2, 1.0, 0.85, RGBColor(0x12, 0x2A, 0x45))
        add_textbox(slide, icon, x + 1.5, 1.22, 1.0, 0.8, font_size=28, align=PP_ALIGN.CENTER)

        add_textbox(slide, role, x, 2.15, card_w, 0.45, font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_rect(slide, x + 0.5, 2.65, 3.0, 0.04, color)

        y = 2.8
        for item in items:
            add_textbox(slide, "✓  " + item, x + 0.2, y, card_w - 0.3, 0.4, font_size=11, color=WHITE)
            y += 0.44

    add_rect(slide, 0, 7.0, 13.33, 0.5, RGBColor(0x0A, 0x14, 0x20))
    add_textbox(slide, "LIVO — Learning Innovation", 0, 7.05, 13.33, 0.35, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    return slide


def add_closing_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, 13.33, 0.08, YELLOW)
    add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)

    add_rect(slide, 1.5, 1.2, 10.3, 5.1, RGBColor(0x0C, 0x1C, 0x30))
    add_rect(slide, 1.5, 1.2, 10.3, 0.1, ACCENT)

    add_textbox(slide, "Terima Kasih", 1.5, 1.6, 10.3, 1.2, font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "LIVO — Learning Innovation", 1.5, 2.7, 10.3, 0.6, font_size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide, 4, 3.5, 5.3, 0.05, YELLOW)

    add_textbox(slide, "Sistem manajemen bimbel yang dirancang untuk\nmemudahkan pengelolaan & meningkatkan kualitas pembelajaran.",
                1.5, 3.7, 10.3, 1.0, font_size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Contact area
    add_rect(slide, 1.5, 4.85, 10.3, 1.0, RGBColor(0x12, 0x2A, 0x45))
    add_textbox(slide, "📍 Srengseng Sawah, Jakarta Selatan     |     📚 Matematika & Bahasa Inggris     |     🎓 TK – SMP",
                1.5, 5.05, 10.3, 0.5, font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)

    return slide


# ─────────────────────────────────────────────────────────────────────────────
def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    add_title_slide(prs)

    # 2. Agenda
    add_content_slide(prs, "Agenda Presentasi", [
        "01 — Profil & Visi LIVO",
        "02 — Program Unggulan (Matematika & Bahasa Inggris)",
        "03 — Pengguna Sistem (Admin, Tutor, Orang Tua)",
        "04 — Fitur Website Publik & Pendaftaran Online",
        "05 — Modul Admin: Dashboard & Manajemen Siswa",
        "06 — Modul Admin: Pembayaran & Verifikasi",
        "07 — Modul Admin: Penjadwalan Terpusat",
        "08 — Modul Admin: Evaluasi & Laporan",
        "09 — Portal Tutor",
        "10 — Alur Sistem End-to-End",
        "11 — Keunggulan Kompetitif",
    ], icon="📋")

    # 3. Section: Profil LIVO
    add_section_divider(prs, "Profil & Visi LIVO", "Bimbingan Belajar Terpadu di Jakarta Selatan")

    add_content_slide(prs, "Tentang LIVO", [
        "## Identitas Lembaga",
        "Nama resmi: LIVO — Learning Innovation",
        "Lokasi: Srengseng Sawah, Jakarta Selatan",
        "Jenjang yang dilayani: TK, SD, SMP",
        "## Visi & Misi",
        "Visi: Menjadi pusat bimbingan belajar terpercaya yang membentuk generasi berpikir kritis & kreatif",
        "Misi: Memberikan pembelajaran terpadu dengan pendekatan personal dan berbasis teknologi",
        "## Nilai Utama",
        "Berpikir Kritis & Kreatif — mengajarkan cara berpikir, bukan sekadar menghafal",
        "Bimbingan Personal — tutor menyesuaikan kecepatan belajar tiap siswa",
        "Dukungan 24/7 — konsultasi tugas tidak terbatas hanya pada jam kelas",
    ], icon="🏫")

    # 4. Program
    add_section_divider(prs, "Program Unggulan", "Dua Program Inti Akademik LIVO")

    add_content_slide(prs, "Program Matematika & Bahasa Inggris", [
        "Program Matematika",
        "Pemahaman konsep dasar hingga lanjutan",
        "Pendekatan berpikir logis dan sistematis",
        "Persiapan ujian sekolah dan kompetisi",
        "Latihan soal dengan pembahasan mendalam",
        "Tutor menyesuaikan kecepatan belajar siswa",
    ], icon="📖", two_col=True, col2_bullets=[
        "Program Bahasa Inggris",
        "Penguasaan kosa kata dan tata bahasa",
        "Metode komunikatif dan kreatif",
        "Latihan reading, writing, dan speaking",
        "Persiapan ujian bahasa sekolah",
        "Pembelajaran kontekstual dan menyenangkan",
    ])

    # 5. Pengguna Sistem
    add_section_divider(prs, "Pengguna Sistem", "Tiga Peran Utama dalam Ekosistem LIVO")
    add_role_slide(prs)

    # 6. Website Publik
    add_section_divider(prs, "Website Publik & Pendaftaran", "Halaman Informasi & Form Registrasi Online")

    add_content_slide(prs, "Fitur Website Publik", [
        "## Halaman Beranda (Landing Page)",
        "Informasi program Matematika & Bahasa Inggris",
        "Statistik: Jenjang TK–SMP, 2 Program, Konsultasi 24/7, Bimbingan 100% Personal",
        "## Keunggulan yang Ditampilkan",
        "Berpikir Kritis & Kreatif — pendekatan holistik",
        "Tutor Berpengalaman — tenaga pengajar terseleksi",
        "Konsultasi Tugas Kapanpun — fleksibel & responsif",
        "## Formulir Pendaftaran Online",
        "Calon siswa mengisi data diri & memilih program",
        "Dukungan kode promo / diskon saat pendaftaran",
        "Admin menerima & memproses data pendaftaran di panel",
    ], icon="🌐")

    # 7. Admin Dashboard
    add_section_divider(prs, "Modul Admin", "Panel Kendali Operasional LIVO")

    add_content_slide(prs, "Dashboard Admin", [
        "## Ringkasan Statistik Real-Time",
        "Total siswa aktif, total sesi belajar, dan pendapatan",
        "Grafik perkembangan pendaftaran bulanan",
        "## Manajemen Pendaftaran",
        "Melihat seluruh pendaftaran masuk dari website",
        "Menyetujui (Approve) atau menolak (Reject) pendaftaran",
        "Membuat akun siswa setelah pendaftaran disetujui",
        "Mencetak kwitansi/bukti pendaftaran (PDF)",
        "## Manajemen Siswa",
        "CRUD lengkap data siswa (tambah, ubah, hapus)",
        "Import data siswa massal via file Excel/template",
        "Melihat detail profil & riwayat sesi siswa",
    ], icon="📊")

    # 8. Pembayaran
    add_section_divider(prs, "Modul Pembayaran", "Verifikasi & Manajemen Transaksi")

    add_content_slide(prs, "Modul Pembayaran — Verifikasi Manual", [
        "## Alur Pembayaran",
        "Orang tua/siswa melakukan transfer bank",
        "Upload foto/dokumen bukti transfer ke sistem",
        "Admin menerima notifikasi konfirmasi pembayaran baru",
        "## Tindakan Admin",
        "Approve (Setujui): sistem otomatis menambah kuota/sesi belajar siswa",
        "Reject (Tolak): pembayaran tidak disetujui, siswa diberitahu",
        "## Fitur Tambahan",
        "Cetak kwitansi pembayaran resmi dalam format PDF",
        "Riwayat seluruh transaksi tersimpan di database",
        "Filter dan pencarian transaksi berdasarkan nama/tanggal",
        "Manajemen paket & promo diskon terintegrasi",
    ], icon="💳")

    # 9. Penjadwalan
    add_section_divider(prs, "Modul Penjadwalan", "Kalender Jadwal Terpusat")

    add_content_slide(prs, "Penjadwalan Terpusat (Admin-Controlled)", [
        "## Master Data Kelas",
        "Pengelolaan paket belajar, nama kelas, & mata pelajaran",
        "Penentuan tarif per sesi belajar",
        "Manajemen silabus per mata pelajaran",
        "## Visualisasi Kalender Interaktif",
        "Kombinasi: Siswa + Tutor + Ruang/Media + Jam",
        "Real-time conflict detection — eliminasi risiko jadwal bentrok",
        "Tampilan mingguan & bulanan",
        "## Manajemen Jadwal",
        "Buat, edit, dan hapus sesi jadwal",
        "Generate jadwal otomatis berdasarkan paket siswa",
        "Update status sesi (aktif / selesai / dibatalkan)",
        "Tutor hanya dapat melihat jadwal (read-only)",
    ], icon="📅")

    # 10. Evaluasi
    add_section_divider(prs, "Modul Evaluasi & Laporan", "Pemantauan Performa Akademik Siswa")

    add_content_slide(prs, "Evaluasi & Laporan Akademik", [
        "## Input oleh Tutor (Per Sesi)",
        "Absensi harian siswa (hadir / tidak hadir / izin)",
        "Nilai skor tugas & kuis harian",
        "Catatan narasi perkembangan akademis & perilaku",
        "## Validasi & Publikasi oleh Admin",
        "Admin memvalidasi data evaluasi dari tutor",
        "Tombol Publish merilis laporan ke halaman orang tua",
        "Laporan tidak terlihat orang tua sebelum dipublikasi",
        "## Output Laporan",
        "Laporan nilai kuis & tugas per siswa",
        "Catatan performa belajar & catatan perilaku",
        "Download ringkasan laporan dalam format PDF",
        "Orang tua dapat memantau perkembangan anak secara langsung",
    ], icon="📝")

    # 11. Portal Tutor
    add_section_divider(prs, "Portal Tutor", "Fitur Khusus untuk Pengajar")

    add_content_slide(prs, "Fitur untuk Tutor", [
        "## Manajemen Biodata",
        "Halaman mandiri untuk melengkapi info profil pribadi",
        "Deklarasi keahlian mata pelajaran yang diampu",
        "Pembaruan nomor kontak aktif",
        "## Kalender Jadwal Mengajar",
        "Visualisasi jadwal dalam format kalender (mingguan/bulanan)",
        "Bersifat read-only — hanya melihat jadwal yang ditetapkan Admin",
        "Tidak dapat mengubah atau membuat jadwal sendiri",
        "## Form Evaluasi Kelas (Input Post-Sesi)",
        "Daftar kehadiran (absensi) siswa per sesi",
        "Input skor nilai tugas & kuis harian",
        "Ruang narasi untuk catatan perkembangan akademis",
        "Catatan perilaku & perkembangan karakter siswa",
    ], icon="👨‍🏫")

    # 12. Alur Sistem
    add_section_divider(prs, "Alur Sistem End-to-End", "Dari Pendaftaran hingga Laporan")

    add_flow_slide(prs, "🔄  Alur Lengkap Sistem LIVO", [
        ("Pendaftaran Online",
         "Calon siswa mengisi formulir di website publik. Data masuk ke panel admin."),
        ("Verifikasi Pendaftaran",
         "Admin meninjau data, menyetujui/menolak, membuat akun siswa."),
        ("Pembayaran & Verifikasi",
         "Siswa transfer + upload bukti. Admin approve → kuota sesi otomatis ditambahkan."),
        ("Penjadwalan Kelas",
         "Admin membuat jadwal berdasarkan paket siswa, tutor, ruang, dan jam tersedia."),
        ("Pelaksanaan Sesi Belajar",
         "Tutor mengajar sesuai jadwal. Silabus & materi sudah tersedia di sistem."),
        ("Input Evaluasi (Tutor)",
         "Tutor mengisi absensi, nilai tugas/kuis, dan catatan setelah sesi selesai."),
        ("Validasi & Publikasi (Admin)",
         "Admin memvalidasi evaluasi, lalu menekan Publish agar terlihat orang tua."),
        ("Monitoring Orang Tua",
         "Orang tua memantau nilai & catatan perkembangan anak secara real-time."),
        ("Siklus Berlanjut",
         "Siswa perpanjang paket → pembayaran baru → jadwal baru → evaluasi baru."),
    ])

    # 13. Keunggulan
    add_section_divider(prs, "Keunggulan Kompetitif", "Mengapa Memilih Sistem LIVO?")

    add_content_slide(prs, "Keunggulan Sistem LIVO", [
        "## Efisiensi Operasional",
        "Semua modul terintegrasi — dari pendaftaran hingga laporan, dalam satu platform",
        "Import data massal via Excel, menghemat waktu entry data",
        "Generate jadwal otomatis berdasarkan paket siswa",
        "## Transparansi & Akuntabilitas",
        "Setiap pembayaran tercatat dengan bukti transfer & kwitansi resmi",
        "Laporan evaluasi terpublikasi — orang tua dapat memantau kapan saja",
        "Riwayat lengkap semua transaksi & sesi tersimpan di database",
        "## Keamanan & Kontrol",
        "Sistem role-based: Admin, Tutor, dan Orang Tua punya akses berbeda",
        "Validasi berlapis — evaluasi harus diverifikasi admin sebelum dipublish",
        "Anti-konflik jadwal: deteksi otomatis benturan waktu/ruang/tutor",
    ], icon="⭐")

    # 14. Closing
    add_closing_slide(prs)

    output = "/Users/mac/Adit Kerjaan/Development/livo/documentation/LIVO_Product_Knowledge.pptx"
    prs.save(output)
    print(f"✅  File berhasil dibuat: {output}")


if __name__ == "__main__":
    build_presentation()
